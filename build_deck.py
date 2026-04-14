from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    gold = RGBColor(196, 164, 132)  # #C4A484
    white = RGBColor(255, 255, 255)
    ink = RGBColor(26, 26, 26)      # #1A1A1A
    parchment = RGBColor(253, 252, 248) # #FDFCF8

    def add_slide(bg_color_name, elements, notes=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Native PowerPoint Background Color
        background = slide.background
        fill = background.fill
        fill.solid()
        if bg_color_name == 'ink':
            fill.fore_color.rgb = ink
        else:
            fill.fore_color.rgb = parchment

        # The "Quiet Line" native visual motif 
        # A faint horizontal line at the top or bottom depending on theme
        if bg_color_name == 'ink':
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2), Inches(8), Inches(14), Inches(8))
            line.line.color.rgb = gold
            line.line.width = Pt(2)
        else:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2), Inches(1), Inches(14), Inches(1))
            line.line.color.rgb = ink
            line.line.width = Pt(4)
        
        for el in elements:
            left, top, w, h = el['rect']
            txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = el['text']
            # Fallback to standard reliable fonts if Playfair Display isn't installed
            p.font.name = el.get('font', 'Georgia')
            p.font.size = Pt(el.get('size', 24))
            p.font.color.rgb = el.get('color', white if bg_color_name == 'ink' else ink)
            p.font.italic = el.get('italic', False)
            p.font.bold = el.get('bold', False)
            if 'align' in el:
                p.alignment = el['align']
            
            for extra in el.get('paragraphs', []):
                p2 = tf.add_paragraph()
                p2.text = extra['text']
                p2.font.name = el.get('font', 'Georgia')
                p2.font.size = Pt(extra.get('size', el.get('size', 24)))
                p2.font.color.rgb = extra.get('color', el.get('color', white if bg_color_name == 'ink' else ink))
                p2.font.italic = extra.get('italic', False)
                p2.font.bold = extra.get('bold', False)
                if 'align' in el:
                    p2.alignment = el['align']

        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes

    # --- SEGMENT 1 ---
    add_slide('ink', [
        {'rect': (2, 3, 12, 2), 'text': 'Professional Drift', 'font': 'Georgia', 'size': 80, 'color': white, 'bold': True, 'align': PP_ALIGN.CENTER},
        {'rect': (2, 5, 12, 1), 'text': 'How Good People Lose Their Way', 'font': 'Georgia', 'size': 36, 'color': gold, 'italic': True, 'align': PP_ALIGN.CENTER}
    ], "Welcome everyone. Thank you for being here. My name is Tristian Walker. Today, we are talking about an invisible force that shapes almost every career.")

    add_slide('ink', [
        {'rect': (2, 4, 12, 1.5), 'text': 'The Day Nobody Meant To Fail', 'font': 'Arial', 'size': 56, 'color': white, 'bold': True, 'align': PP_ALIGN.CENTER},
        {'rect': (2, 5.5, 12, 1.5), 'text': 'April 14, 1994. 26 Lives. 0 Intentional Mistakes.', 'font': 'Georgia', 'size': 32, 'color': gold, 'italic': True, 'align': PP_ALIGN.CENTER}
    ], "Let’s start with a story. April 14, 1994. The Black Hawk Shootdown over Iraq. Two U.S. Black Hawk helicopters were shot down by friendly fire. Nobody made one catastrophic decision. There were dozens of small adjustments, tiny deviations, little accommodations.")

    add_slide('ink', [
        {'rect': (2, 3.5, 12, 2.5), 'text': '“Practical Drift is the slow, steady uncoupling of local practice from written procedure.”', 'font': 'Georgia', 'size': 44, 'color': white, 'italic': True, 'align': PP_ALIGN.CENTER},
        {'rect': (2, 6.5, 12, 1), 'text': '— Scott A. Snook, Friendly Fire', 'font': 'Arial', 'size': 24, 'color': gold, 'align': PP_ALIGN.CENTER}
    ], "Researcher Scott A. Snook calls this 'Practical Drift' — the 'slow, steady uncoupling of local practice from written procedure.' What happened in the skies over Iraq happens in conference rooms, cubicles, and corner offices every single day. Not with the same finality — but with the same quiet, devastating certainty.")

    add_slide('ink', [
        {'rect': (2, 3, 12, 2), 'text': 'The Quiet Line He Never Saw', 'font': 'Georgia', 'size': 56, 'color': gold, 'bold': True, 'align': PP_ALIGN.CENTER},
        {'rect': (2, 5, 12, 1), 'text': 'The line did not arrive as a warning. \nIt arrived as routine.', 'font': 'Georgia', 'size': 36, 'color': white, 'italic': True, 'align': PP_ALIGN.CENTER}
    ], "The line did not arrive as a warning. It arrived as routine. Terrence traveled constantly, and his life became a sequence of early flights, rental counters, and hotel check-ins.\n\nIn the private silence of hotel rooms, when the hallway noise faded and the television offered only shapes without sound, Terrence found himself drinking because he did not want to feel the sameness. Not sadness, exactly. The flatness that accumulates when your days have lost their edges.\n\nThe line never looked like a boundary; it looked like normal. And Terrence never saw the moment he crossed it. He was moving, but no longer toward anything he had chosen.")

    # --- SEGMENT 2 ---
    add_slide('ink', [
        {'rect': (2, 3, 12, 2), 'text': 'Drift Is Not a Detour.', 'font': 'Georgia', 'size': 64, 'color': white, 'align': PP_ALIGN.CENTER, 'italic': True},
        {'rect': (2, 4.5, 12, 2), 'text': 'It’s a Direction.', 'font': 'Georgia', 'size': 64, 'color': gold, 'align': PP_ALIGN.CENTER, 'italic': True}
    ], "Professional Drift is what happens when you are moving — but not toward anything you chose.")

    add_slide('ink', [
        {'rect': (1, 1.5, 14, 1.5), 'text': 'The Taxonomy of Drift', 'font': 'Georgia', 'size': 48, 'color': gold, 'bold': True, 'align': PP_ALIGN.CENTER},
        {'rect': (1, 3.5, 4.5, 4), 'text': 'Reactive Drift', 'font': 'Georgia', 'size': 36, 'color': gold, 'bold': True, 'align': PP_ALIGN.CENTER, 'paragraphs': [{'text': '\nWhen you stop driving and start responding. \n(Hyatt & Harkavy)', 'size': 24, 'color': white}]},
        {'rect': (5.5, 3.5, 5, 4), 'text': 'Success Trap', 'font': 'Georgia', 'size': 36, 'color': gold, 'bold': True, 'align': PP_ALIGN.CENTER, 'paragraphs': [{'text': '\nThe applause gets louder the further you walk from yourself.\n(Junko Okada)', 'size': 24, 'color': white}]},
        {'rect': (10.5, 3.5, 4.5, 4), 'text': 'Systemic Drift', 'font': 'Georgia', 'size': 36, 'color': gold, 'bold': True, 'align': PP_ALIGN.CENTER, 'paragraphs': [{'text': '\nYour environment changes; your identity doesn’t update.', 'size': 24, 'color': white}]}
    ], "We see three types of drift. Reactive Drift, the Success Trap, and Systemic Drift. Motion is when you are busy. Movement is when you are somewhere new at the end of it.")

    add_slide('ink', [
        {'rect': (2, 3.5, 12, 3), 'text': 'Raise your hand if you’ve had the most productive week of your career—and still felt like something was wrong.', 'font': 'Georgia', 'size': 48, 'color': white, 'italic': True, 'align': PP_ALIGN.CENTER}
    ], "[Interactive Moment] Ask the audience: Raise your hand if you've had the most productive week of your career and still felt like something was wrong.")

    # --- SEGMENT 3 ---
    add_slide('parchment', [
        {'rect': (2, 1.5, 12, 1.5), 'text': 'You Can’t Correct What You Can’t Name', 'font': 'Georgia', 'size': 48, 'color': ink, 'bold': True, 'align': PP_ALIGN.CENTER},
        {'rect': (2, 3, 12, 5), 'text': '1. The Shrinking Why - You can\'t explain why you\'re doing what you\'re doing anymore.\n\n2. The Expertise Plateau - You haven\'t learned something genuinely new in 6+ months.\n\n3. The Validation Loop - External success has replaced internal purpose.\n\n4. The Reactive Calendar - You serve others\' priorities exclusively.\n\n5. The Identity Gap - Who you are at work no longer matches who you believe you are.', 'font': 'Arial', 'size': 28, 'color': ink, 'align': PP_ALIGN.LEFT}
    ], "Here are the 5 signals of drift. The Shrinking Why, The Expertise Plateau, The Validation Loop, The Reactive Calendar, and The Identity Gap.")

    add_slide('parchment', [
        {'rect': (2, 2.5, 12, 6), 'text': '“He is the file now. The reduction he had always resisted in others has arrived for him, clinical and complete...”\n\n“This is the Identity Gap: the quiet, terrifying distance between the man who leads and the man being breathed for by a machine.”', 'font': 'Georgia', 'size': 36, 'color': ink, 'italic': True, 'align': PP_ALIGN.CENTER}
    ], "[READING MOMENT]\nHe had been in enough professional rooms to recognize the moment a name becomes a file number, the moment a person stops being seen and starts being processed. He tries to answer—to say 'I am here'—but his tongue feels enormous and useless, a foreign object lodged behind his teeth...\n\nWhat brought him here had not arrived without warning. It had arrived as habit, the way most things do... the quiet, terrifying distance between the man who leads and the man being breathed for by a machine.")

    add_slide('parchment', [
        {'rect': (2, 2.5, 12, 4), 'text': 'Are you on a Bridge or a Detour?', 'font': 'Georgia', 'size': 56, 'color': ink, 'bold': True, 'align': PP_ALIGN.CENTER, 'paragraphs': [{'text': '\n\nA Bridge takes you somewhere new.\nA Detour feels like progress but loops you back to where you started.', 'size': 32, 'color': ink, 'italic': True}]}
    ], "Ask yourself the Faststream diagnostic: Where will I be in 5 years if I keep doing exactly this? Am I being stretched or just busy?")

    # --- SEGMENT 4 ---
    add_slide('parchment', [
        {'rect': (2, 3, 12, 2), 'text': 'Drift Is Not Your Destiny.', 'font': 'Georgia', 'size': 64, 'color': ink, 'bold': True, 'align': PP_ALIGN.CENTER},
        {'rect': (2, 5, 12, 2), 'text': 'It is Data.', 'font': 'Georgia', 'size': 64, 'color': ink, 'italic': True, 'align': PP_ALIGN.CENTER}
    ], "Shift the narrative: Drift is not a character flaw. It is feedback from your life. The fact that you notice the drift means the part of you that knows where you're supposed to be is still awake.")

    add_slide('parchment', [
        {'rect': (2, 3.5, 12, 3), 'text': 'The Values Compass', 'font': 'Georgia', 'size': 56, 'color': ink, 'bold': True, 'align': PP_ALIGN.CENTER, 'paragraphs': [{'text': '\nRigid long-term goal-setting breaks when terrain changes. A compass gives you direction, not destination.', 'size': 36, 'italic': True, 'color': ink}]}
    ], "Replace the Destination Map with a Values Compass. Jen Green talks about Evolutionary Scope: Drift is dangerous when unconscious, but powerful when intentional.")

    add_slide('parchment', [
        {'rect': (2, 3.5, 12, 3), 'text': 'Treat your career like a prototype, not a contract.', 'font': 'Georgia', 'size': 52, 'color': ink, 'italic': True, 'align': PP_ALIGN.CENTER, 'paragraphs': [{'text': '\nYou are allowed to iterate. Iteration is not failure.', 'size': 36, 'color': ink}]}
    ], "In Designing Your Life, Burnett & Evans teach us to treat our careers like prototypes. Beware 'The Gravity Problem' — some forces you can't move, but you can navigate around them.")

    add_slide('parchment', [
        {'rect': (2, 1.5, 12, 1.5), 'text': 'In which domain has the drift been the loudest?', 'font': 'Georgia', 'size': 48, 'color': ink, 'bold': True, 'align': PP_ALIGN.CENTER},
        {'rect': (2, 3.5, 12, 4), 'text': '1. Career        2. Relationships        3. Health\n\n4. Finances        5. Personal Development\n\n6. Spirituality        7. Recreation', 'font': 'Arial', 'size': 32, 'color': ink, 'align': PP_ALIGN.CENTER}
    ], "Living Forward's 7 Domains (Hyatt & Harkavy): Career. Relationships. Health. Finances. Personal development. Spirituality. Recreation. The answer tells you where to start.")

    # --- SEGMENT 5 ---
    add_slide('parchment', [
        {'rect': (2, 3.5, 12, 3), 'text': '“Every single moment of drift had a preceding moment of awareness. A flicker. A hesitation.”', 'font': 'Georgia', 'size': 48, 'color': ink, 'italic': True, 'align': PP_ALIGN.CENTER}
    ], "Call back to the opening. In 1994, it took multiple small drifts... But Snook also found: every single moment of drift had a preceding moment of awareness. A flicker. A hesitation. People knew, and then they adapted anyway.")

    add_slide('ink', [
        {'rect': (1, 2.5, 14, 5), 'text': '“The path does not disappear. It begins again the moment you decide to walk it. This is the manifesto of the flicker: that we are not defined by the drift, but by the quiet, disciplined return to the room.”', 'font': 'Georgia', 'size': 44, 'color': gold, 'italic': True, 'align': PP_ALIGN.CENTER}
    ], "The 'Quiet Line' is the decision to reclaim that flicker. It is the realization that drift does not feel like failure—that is what makes it dangerous... The path begins again the moment you decide to walk it. Leave with this: What would I be doing professionally if I trusted myself completely?")

    output_path = "/Users/tristianwalker/Hospitality Lecture Series /professional-drift/Quiet_Line_V2_Final.pptx"
    prs.save(output_path)
    print(f"Deck rebuilt dynamically at: {output_path}")

if __name__ == "__main__":
    create_deck()
